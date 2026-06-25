import io
from pathlib import Path
from typing import Any, Literal

import pandas as pd
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle
from pptx import Presentation
from pptx.util import Pt

from app.services.dashboard_service import generate_dashboard
from app.services.forecast_engine import run_forecast_leaderboard
from app.services.health_score_service import compute_health_score
from app.services.insight_engine import generate_insights
from app.services.report_engine import generate_report
from app.services.storytelling_service import generate_story

_LOGO_CANDIDATES = [
    Path(__file__).resolve().parent.parent.parent.parent / "assets" / "logo.png",
    Path(__file__).resolve().parent.parent.parent / "assets" / "logo.png",
]


async def generate_full_report(
    df: pd.DataFrame,
    filename: str,
    fmt: Literal["markdown", "pdf", "pptx"] = "markdown",
) -> dict[str, Any]:
    base = await generate_report(df, filename)
    health = compute_health_score(df)
    insights = generate_insights(df)
    story = await generate_story(df)
    dashboard = generate_dashboard(df)

    forecast_outlook = "Forecast not available for this dataset."
    forecast_data: dict[str, Any] | None = None
    try:
        fc = run_forecast_leaderboard(df, forecast_horizon=30)
        if fc.get("available") and fc.get("forecast"):
            forecast_data = fc
            exec_sum = fc.get("executive_summary", {})
            next_val = fc["forecast"][-1]["value"]
            forecast_outlook = (
                f"Best model: {fc['best_model']['model_name']}. "
                f"End-of-horizon forecast for {fc['target_column']}: {next_val:,.2f}. "
                f"{exec_sum.get('ai_commentary', fc.get('explanation', ''))}"
            )
    except ValueError:
        pass

    model_results = "Model Arena results not available for this dataset."
    try:
        from app.services.model_arena_service import run_model_arena

        arena = run_model_arena(df)
        model_results = arena.get("performance_summary", "")
    except ValueError:
        pass

    key_findings = [i["description"] for i in insights[:5]]
    risks = [i["description"] for i in insights if i["severity"] == "high"][:5]
    if not risks and health["overall_score"] < 60:
        risks = [f"Data quality score is {health['overall_score']}/100."]
    opportunities = [
        i["description"] for i in insights
        if i["type"] in ("growth", "category_performance", "trend")
    ][:5]
    recommendations = health.get("recommended_fixes", [])[:5]
    recommendations.extend(story["what_to_do_next"].split(". ")[:2])
    recommendations = [r.strip() for r in recommendations if r.strip()][:8]

    executive_summary = (
        f"Analysis of '{filename}' ({health['rows']} rows, {health['columns']} columns). "
        f"Health score: {health['overall_score']}/100. "
        f"{story['what_happened']}"
    )

    result: dict[str, Any] = {
        "markdown": base["markdown"],
        "executive_summary": executive_summary,
        "key_findings": key_findings,
        "risks": risks,
        "opportunities": opportunities,
        "recommendations": recommendations,
        "forecast_outlook": forecast_outlook,
        "llm_enhanced": base.get("llm_enhanced", False),
        "format": fmt,
    }

    if fmt == "pdf":
        result["file_bytes"] = _build_pdf(
            filename,
            executive_summary,
            health,
            key_findings,
            story.get("why_it_happened", ""),
            forecast_outlook,
            forecast_data,
            model_results,
            risks,
            opportunities,
            recommendations,
            dashboard.get("kpis", []),
        )
    elif fmt == "pptx":
        result["file_bytes"] = _build_pptx(
            filename, executive_summary, key_findings, risks, opportunities, recommendations
        )

    return result


def _logo_path() -> Path | None:
    for path in _LOGO_CANDIDATES:
        if path.is_file():
            return path
    return None


def _build_pdf(
    filename: str,
    summary: str,
    health: dict[str, Any],
    findings: list[str],
    root_cause: str,
    forecast_outlook: str,
    forecast_data: dict[str, Any] | None,
    model_results: str,
    risks: list[str],
    opportunities: list[str],
    recommendations: list[str],
    kpis: list[dict[str, Any]],
) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    cover_title = ParagraphStyle(
        "CoverTitle",
        parent=styles["Title"],
        fontSize=28,
        spaceAfter=12,
        textColor=colors.HexColor("#18181b"),
    )
    cover_sub = ParagraphStyle(
        "CoverSub",
        parent=styles["Normal"],
        fontSize=12,
        textColor=colors.HexColor("#71717a"),
        spaceAfter=6,
    )
    section = ParagraphStyle(
        "Section",
        parent=styles["Heading2"],
        fontSize=14,
        spaceBefore=16,
        spaceAfter=8,
        textColor=colors.HexColor("#27272a"),
    )
    story: list[Any] = []

    logo = _logo_path()
    if logo:
        story.append(Image(str(logo), width=1.0 * inch, height=1.0 * inch))
        story.append(Spacer(1, 24))

    story.append(Paragraph("Prisma AI", cover_sub))
    story.append(Paragraph("Executive Analytics Report", cover_title))
    story.append(Spacer(1, 8))
    story.append(Paragraph(f"<b>Dataset:</b> {filename}", cover_sub))
    story.append(Paragraph(
        f"<b>Scope:</b> {health['rows']:,} rows · {health['columns']} columns · "
        f"Health score {health['overall_score']}/100",
        cover_sub,
    ))
    story.append(Spacer(1, 36))
    story.append(Paragraph("<i>Confidential — Prepared by Prisma AI Analytics Platform</i>", cover_sub))
    story.append(PageBreak())

    story.append(Paragraph("Dataset Summary", section))
    summary_table = Table(
        [
            ["Metric", "Value"],
            ["Rows", f"{health['rows']:,}"],
            ["Columns", str(health["columns"])],
            ["Health Score", f"{health['overall_score']}/100"],
            ["Completeness", f"{health['sub_scores']['completeness']:.0f}"],
        ],
        colWidths=[2.5 * inch, 3.5 * inch],
    )
    summary_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f4f4f5")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#52525b")),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#e4e4e7")),
    ]))
    story.append(summary_table)
    story.append(Spacer(1, 12))
    story.append(Paragraph("<b>Executive Summary</b>", styles["Heading3"]))
    story.append(Paragraph(summary, styles["Normal"]))

    if kpis:
        story.append(Spacer(1, 12))
        story.append(Paragraph("<b>Key Metrics</b>", styles["Heading3"]))
        for k in kpis[:4]:
            story.append(
                Paragraph(
                    f"• {k.get('label', k.get('column'))}: {k.get('value', 0):,.2f}",
                    styles["Normal"],
                )
            )

    sections = [
        ("Key Findings", findings),
        ("Root Cause Analysis", [root_cause] if root_cause else []),
        ("Forecasts", [forecast_outlook]),
        ("Model Results", [model_results]),
        ("Risks", risks),
        ("Opportunities", opportunities),
        ("Recommendations", recommendations),
    ]
    for title, items in sections:
        story.append(Paragraph(title, section))
        for item in items or ["None identified for this dataset."]:
            story.append(Paragraph(f"• {item}", styles["Normal"]))
        story.append(Spacer(1, 4))

    if forecast_data and forecast_data.get("executive_summary"):
        es = forecast_data["executive_summary"]
        story.append(Paragraph("Forecast Detail", section))
        forecast_rows = [
            ["Indicator", "Value"],
            ["Current Trend", f"{es.get('current_trend', '—')} ({es.get('trend_pct', 0):+.1f}%)"],
            ["Projected Change", f"{es.get('projected_change_pct', 0):+.1f}%"],
            ["Confidence", es.get("confidence_level", "—")],
            ["Best Case", f"{es.get('best_case', 0):,.2f}"],
            ["Worst Case", f"{es.get('worst_case', 0):,.2f}"],
        ]
        ft = Table(forecast_rows, colWidths=[2.5 * inch, 3.5 * inch])
        ft.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#f4f4f5")),
            ("FONTSIZE", (0, 0), (-1, -1), 10),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#e4e4e7")),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
        ]))
        story.append(ft)

    doc.build(story)
    return buffer.getvalue()


def _build_pptx(
    filename: str,
    summary: str,
    findings: list[str],
    risks: list[str],
    opportunities: list[str],
    recommendations: list[str],
) -> bytes:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Executive Report"
    title_slide.placeholders[1].text = filename

    sections = [
        ("Executive Summary", [summary]),
        ("Key Findings", findings),
        ("Risks", risks),
        ("Opportunities", opportunities),
        ("Recommendations", recommendations),
    ]

    for section_title, bullets in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = section_title
        body = slide.placeholders[1].text_frame
        body.clear()
        for i, bullet in enumerate(bullets or ["None identified."]):
            p = body.paragraphs[0] if i == 0 else body.add_paragraph()
            p.text = bullet
            p.font.size = Pt(14)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
